from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.shapes.picture import Picture
from pptx.shapes.group import GroupShape
from .utils import get_img_dir
from ..formulas import OMML_NS, FormulaRef, ExtractionResult

_OMATH_PARA = f"{{{OMML_NS}}}oMathPara"
_OMATH = f"{{{OMML_NS}}}oMath"


def extract_pptx(file_path: Path, temp_dir: Path) -> ExtractionResult:
    prs = Presentation(str(file_path))
    img_dir = get_img_dir(temp_dir, file_path)

    lines = []
    img_count = 0
    formulas: list[FormulaRef] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        lines.append(f"\n## Slide {slide_num}\n")
        for shape in _iter_shapes(slide.shapes):
            if isinstance(shape, Picture):
                try:
                    ext = _img_ext(shape.image.content_type)
                    img_path = img_dir / f"slide{slide_num}_img{img_count}{ext}"
                    img_path.write_bytes(shape.image.blob)
                    lines.append(f"\n![image]({img_path.resolve()})\n")
                    img_count += 1
                except (ValueError, AttributeError):
                    pass
            elif shape.has_text_frame:
                tf_elem = shape.text_frame._txBody
                shape_formulas = _detect_shape_formulas(
                    tf_elem, slide_num, shape.text_frame.text.strip(), formulas
                )

                if shape_formulas:
                    # Build text with formula markers
                    text = _build_shape_text_with_markers(
                        shape, shape_formulas, len(formulas) - len(shape_formulas)
                    )
                    if text:
                        lines.append(text)
                else:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)
            elif shape.has_table:
                lines.append(_table_to_md(shape.table))

        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"\n> **Notes:** {notes}\n")

    return ExtractionResult(
        markdown="\n".join(lines),
        formulas=formulas,
    )


def _detect_shape_formulas(
    tf_elem, slide_num: int, context: str, formulas_list: list
) -> list[FormulaRef]:
    """Detect OMML formulas in a text frame element."""
    if tf_elem is None:
        return []

    found = []

    # Display formulas
    omath_paras = tf_elem.findall(f".//{_OMATH_PARA}")
    for omp in omath_paras:
        xml_str = etree.tostring(omp, encoding="unicode")
        ref = FormulaRef(
            omml_xml=xml_str,
            is_inline=False,
            context_text=f"[Slide {slide_num}] {context[:200]}",
        )
        formulas_list.append(ref)
        found.append(ref)

    # Inline formulas
    all_omath = tf_elem.findall(f".//{_OMATH}")
    for om in all_omath:
        parent = om.getparent()
        if parent is not None and parent.tag == _OMATH_PARA:
            continue
        xml_str = etree.tostring(om, encoding="unicode")
        ref = FormulaRef(
            omml_xml=xml_str,
            is_inline=True,
            context_text=f"[Slide {slide_num}] {context[:200]}",
        )
        formulas_list.append(ref)
        found.append(ref)

    return found


def _build_shape_text_with_markers(
    shape, shape_formulas: list[FormulaRef], start_idx: int
) -> str:
    """Build text frame content with formula markers."""
    parts = []
    formula_idx = 0

    for para in shape.text_frame.paragraphs:
        para_text = para.text.strip()
        if para_text:
            parts.append(para_text)

    # If all formulas are display mode, insert markers on their own lines
    result_parts = []
    formula_offset = 0
    for ref in shape_formulas:
        if not ref.is_inline:
            result_parts.append(ref.marker(start_idx + formula_offset))
        formula_offset += 1

    # If we have inline formulas mixed with text, just append markers after text
    inline_formulas = [f for f in shape_formulas if f.is_inline]
    if inline_formulas and parts:
        # Simple approach: append inline formula markers to the text content
        text = "\n".join(parts)
        for i, ref in enumerate(inline_formulas):
            idx = start_idx + shape_formulas.index(ref)
            text += f" {ref.marker(idx)}"
        result_parts.insert(0, text)
    elif parts and not result_parts:
        result_parts = parts

    return "\n".join(result_parts)


def _iter_shapes(shapes):
    for shape in shapes:
        if isinstance(shape, GroupShape):
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _table_to_md(table) -> str:
    rows = []
    for row in table.rows:
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        rows.append(cells)
    if not rows:
        return ""
    header = "| " + " | ".join(rows[0]) + " |"
    sep = "| " + " | ".join("---" for _ in rows[0]) + " |"
    body = "\n".join("| " + " | ".join(r) + " |" for r in rows[1:])
    return f"\n{header}\n{sep}\n{body}\n"


def _img_ext(content_type: str) -> str:
    return {"image/png": ".png", "image/jpeg": ".jpg", "image/gif": ".gif",
            "image/bmp": ".bmp", "image/tiff": ".tiff"}.get(content_type, ".png")
