"""Detect OMML math formulas in PPTX documents."""

import logging
from lxml import etree
from pptx import Presentation
from pptx.shapes.group import GroupShape

from . import OMML_NS, FormulaRef

log = logging.getLogger("knowledge_extractor")

_OMATH_PARA = f"{{{OMML_NS}}}oMathPara"
_OMATH = f"{{{OMML_NS}}}oMath"

# The PPTX math namespace uses the same URI as DOCX (a14:m in the XML,
# but resolved to the same OMML namespace)
_A14_NS = "http://schemas.microsoft.com/office/drawing/2010/main"


def detect_formulas(prs: Presentation) -> list[FormulaRef]:
    """Walk a PPTX presentation and detect all OMML formula elements.

    Returns a list of FormulaRef objects. Formulas in PPTX are typically
    display mode unless found inline within a text run.
    """
    formulas = []

    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in _iter_shapes(slide.shapes):
            if not shape.has_text_frame:
                continue

            tf_elem = shape.text_frame._txBody
            if tf_elem is None:
                continue

            # Get text context from the shape
            context = ""
            if shape.has_text_frame:
                context = shape.text_frame.text.strip()[:300]

            # Search for oMathPara (display mode)
            omath_paras = tf_elem.findall(f".//{_OMATH_PARA}")
            for omp in omath_paras:
                xml_str = etree.tostring(omp, encoding="unicode")
                formulas.append(FormulaRef(
                    omml_xml=xml_str,
                    is_inline=False,
                    context_text=f"[Slide {slide_num}] {context}",
                ))

            # Search for bare oMath (inline mode)
            all_omath = tf_elem.findall(f".//{_OMATH}")
            for om in all_omath:
                parent = om.getparent()
                if parent is not None and parent.tag == _OMATH_PARA:
                    continue
                xml_str = etree.tostring(om, encoding="unicode")
                formulas.append(FormulaRef(
                    omml_xml=xml_str,
                    is_inline=True,
                    context_text=f"[Slide {slide_num}] {context}",
                ))

    log.debug(f"PPTX formula detection: found {len(formulas)} formulas")
    return formulas


def _iter_shapes(shapes):
    """Recursively iterate all shapes, including inside groups."""
    for shape in shapes:
        if isinstance(shape, GroupShape):
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape
